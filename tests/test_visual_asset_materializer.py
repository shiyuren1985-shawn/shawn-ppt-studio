from __future__ import annotations

import importlib.util
import json
from pathlib import Path
import sys
import tempfile
import unittest


SCRIPTS = Path(__file__).resolve().parents[1] / "scripts"
sys.path.insert(0, str(SCRIPTS))
SPEC = importlib.util.spec_from_file_location(
    "visual_asset_materializer", SCRIPTS / "materialize_visual_asset.py"
)
assert SPEC and SPEC.loader
materializer = importlib.util.module_from_spec(SPEC)
SPEC.loader.exec_module(materializer)


def write_minimal_pdf(path: Path) -> None:
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 200 100] /Resources << >> /Contents 4 0 R >>",
        b"<< /Length 0 >>\nstream\n\nendstream",
    ]
    payload = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, value in enumerate(objects, start=1):
        offsets.append(len(payload))
        payload.extend(f"{index} 0 obj\n".encode())
        payload.extend(value)
        payload.extend(b"\nendobj\n")
    xref_offset = len(payload)
    payload.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    payload.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        payload.extend(f"{offset:010d} 00000 n \n".encode())
    payload.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode()
    )
    path.write_bytes(payload)


class VisualAssetMaterializerTests(unittest.TestCase):
    def test_pdf_page_is_materialized_with_reusable_receipt(self) -> None:
        with tempfile.TemporaryDirectory(prefix="visual_asset_pdf_") as temp:
            root = Path(temp).resolve()
            source = root / "source.pdf"
            output = root / "assets"
            write_minimal_pdf(source)
            spec = {
                "source": str(source),
                "role": "source_page",
                "locator": {"page": 1},
            }
            try:
                first = materializer.materialize_visual_asset(spec, output)
            except SystemExit as exc:
                if "缺少视觉资产转换器" in str(exc):
                    self.skipTest("pdftoppm is not installed")
                raise
            second = materializer.materialize_visual_asset(
                {**spec, "role": "case_evidence"}, output
            )
            rendered = Path(first["output_path"])
            self.assertEqual(rendered.read_bytes()[:8], b"\x89PNG\r\n\x1a\n")
            self.assertEqual(first["output_path"], second["output_path"])
            self.assertEqual(second["role"], "case_evidence")
            self.assertEqual(first["receipt_path"], second["receipt_path"])

    def test_pptx_slide_is_materialized_through_shared_office_adapter(self) -> None:
        try:
            from pptx import Presentation
        except ImportError:
            self.skipTest("python-pptx is not installed")
        with tempfile.TemporaryDirectory(prefix="visual_asset_pptx_") as temp:
            root = Path(temp).resolve()
            source = root / "source.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "Shared visual materializer"
            presentation.save(source)
            try:
                result = materializer.materialize_visual_asset(
                    {
                        "source": str(source),
                        "role": "source_slide",
                        "locator": {"slide": 1},
                    },
                    root / "assets",
                )
            except SystemExit as exc:
                if "缺少视觉资产转换器" in str(exc):
                    self.skipTest("Office or PDF renderer is not installed")
                raise
            self.assertEqual(result["source_kind"], "presentation_slide")
            self.assertEqual(Path(result["output_path"]).suffix, ".png")
            self.assertEqual(result["renderer"], "soffice+pdftoppm")

    def test_local_html_is_captured_as_a_web_viewport(self) -> None:
        with tempfile.TemporaryDirectory(prefix="visual_asset_web_") as temp:
            root = Path(temp).resolve()
            source = root / "source.html"
            source.write_text(
                "<!doctype html><style>body{margin:0;background:#123;color:white}"
                "h1{font:64px sans-serif}</style><h1>Visual source</h1>",
                encoding="utf-8",
            )
            try:
                result = materializer.materialize_visual_asset(
                    {
                        "source": str(source),
                        "role": "web_reference",
                        "locator": {
                            "viewport": {"width": 800, "height": 450},
                            "wait_ms": 100,
                        },
                    },
                    root / "assets",
                )
            except SystemExit as exc:
                if "缺少视觉资产转换器" in str(exc):
                    self.skipTest("Chrome is not installed")
                raise
            self.assertEqual(result["source_kind"], "web_viewport")
            self.assertEqual(result["locator"]["viewport"], {"width": 800, "height": 450})
            receipt = json.loads(
                Path(result["receipt_path"]).read_text(encoding="utf-8")
            )
            self.assertNotIn("role", receipt)

    def test_unknown_source_type_fails_without_a_fake_raster(self) -> None:
        with tempfile.TemporaryDirectory(prefix="visual_asset_unknown_") as temp:
            root = Path(temp).resolve()
            source = root / "source.xyz"
            source.write_text("unknown", encoding="utf-8")
            with self.assertRaises(SystemExit) as context:
                materializer.materialize_visual_asset(
                    {"source": str(source), "role": "evidence"},
                    root / "assets",
                )
            self.assertIn("尚不支持", str(context.exception))
            self.assertEqual(list((root / "assets").glob("*.png")), [])


if __name__ == "__main__":
    unittest.main()
